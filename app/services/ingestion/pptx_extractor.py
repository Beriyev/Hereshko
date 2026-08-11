from pptx import Presentation
from app.core.normalization import Document, SourceType
import uuid
from pathlib import Path
from datetime import datetime, timezone
from app.core.exceptions import IngestionError

def extract_pptx(file_path: Path, notebook_id: str) -> Document:
    texts = []
    boundaries = []
    offset = 0
    slide_number = 1
    slides_count = 0
    content = ""

    try:
        presentation = Presentation(str(file_path))
        slides_count = len(presentation.slides)
        for slide in presentation.slides:
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = getattr(shape, 'text', "")
                    if text:
                        slide_text += text + "\n"
            if slide_text.strip():
                slide_text = slide_text.strip()
                text_size = len(slide_text)
                boundaries.append({
                    "slide_number": slide_number,
                    "start": offset,
                    "end": offset + text_size
                })
                texts.append(slide_text)
                offset += text_size+1
            slide_number += 1
    except Exception as e:
        raise IngestionError(f"Error occurred while extracting PPTX file: {str(e)}")
    content = "\n".join(texts) + "\n" if texts else ""
    if not content.strip():
        raise IngestionError("No extractable text found in PPTX (may be an empty presentation)")

    return Document(
       document_id = str(uuid.uuid4()),
        notebook_id = notebook_id,
        content = content,
        source_type = SourceType.PPTX,
        source_identifier = file_path.name,
        title = file_path.stem,
        ingested_at = datetime.now(timezone.utc),
        raw_metadata = {
            "original_filename": file_path.name,
            "file_size_bytes": file_path.stat().st_size,
            "slides_count": slides_count,
            "boundaries": boundaries
        }
    )